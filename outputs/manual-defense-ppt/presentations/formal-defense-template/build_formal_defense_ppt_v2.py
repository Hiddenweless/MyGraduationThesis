from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt, Inches


ROOT = Path("/Users/wtl/Desktop/我的毕业论文 ")
WORK = ROOT / "outputs/manual-defense-ppt/presentations/formal-defense-template"
ASSETS = WORK / "assets"
OUT_DIR = WORK / "output"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUT_DIR / "王腾龙-硕士论文正式答辩PPT-优化版.pptx"
OUTLINE_OUT = OUT_DIR / "王腾龙-硕士论文正式答辩PPT-优化版大纲.md"
REVIEW_OUT = OUT_DIR / "二轮对比反省与优化说明.md"

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

W, H = prs.slide_width, prs.slide_height
BLUE = RGBColor(2, 64, 154)
BLUE2 = RGBColor(33, 94, 167)
BLUE_LIGHT = RGBColor(226, 236, 250)
ORANGE = RGBColor(237, 125, 49)
GRAY = RGBColor(89, 89, 89)
BLACK = RGBColor(35, 35, 35)
WHITE = RGBColor(255, 255, 255)
PALE = RGBColor(246, 249, 253)
FONT = "微软雅黑"

slides_meta = []


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def textbox(slide, x, y, w, h, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    sh = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return sh


def bullets(slide, x, y, w, h, items, size=16, color=BLACK):
    sh = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        for r in p.runs:
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return sh


def header(slide, title, page=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Cm(2.35))
    fill(bar, BLUE)
    for i, h in enumerate([2.35, 1.55, 0.85]):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.25 + i * 0.18), 0, Cm(0.03), Cm(h))
        fill(line, WHITE)
    textbox(slide, 1.72, 0.42, 22.0, 1.3, title, size=24, color=WHITE, bold=True)
    if page is not None:
        textbox(slide, 10.5, 18.0, 4.0, 0.5, str(page), size=12, color=GRAY, align=PP_ALIGN.CENTER)


def add_img(slide, name, x, y, w, h):
    p = ASSETS / name
    if p.exists():
        slide.shapes.add_picture(str(p), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    else:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
        fill(sh, PALE)
        sh.line.color.rgb = GRAY
        textbox(slide, x + 0.5, y + h / 2 - 0.3, w - 1.0, 0.6, f"需替换：{name}", size=12, color=GRAY, align=PP_ALIGN.CENTER)


def box(slide, x, y, w, h, text="", color=PALE, line=BLUE2):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.color.rgb = line
    if text:
        textbox(slide, x + 0.25, y + 0.15, w - 0.5, h - 0.3, text, size=15, color=BLACK)
    return sh


def tag(slide, x, y, text, w=3.4, color=BLUE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(0.9))
    fill(sh, color)
    textbox(slide, x, y + 0.12, w, 0.5, text, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return sh


def table(slide, x, y, data, colw, fs=10):
    rows, cols = len(data), len(data[0])
    shp = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(sum(colw)), Cm(rows * 0.78))
    tbl = shp.table
    for c, w in enumerate(colw):
        tbl.columns[c].width = Cm(w)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Cm(0.06)
            cell.margin_right = Cm(0.06)
            cell.margin_top = Cm(0.03)
            cell.margin_bottom = Cm(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE if r == 0 else (RGBColor(241, 245, 251) if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for rr in p.runs:
                    rr.font.name = FONT
                    rr.font.size = Pt(fs)
                    rr.font.bold = r == 0
                    rr.font.color.rgb = WHITE if r == 0 else BLACK
    return tbl


def section(page, title, active):
    s = blank()
    header(s, "提纲", page)
    items = [("1", "研究背景、研究现状"), ("2", "研究目标、研究内容"), ("3", "技术路线、系统实现"), ("4", "实验验证、总结展望")]
    y = 4.1
    for n, t in items:
        sh = slide_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(6.2), Cm(y), Cm(13.0), Cm(1.6))
        sh.fill.solid()
        sh.fill.fore_color.rgb = BLUE if n == str(active) else BLUE_LIGHT
        sh.line.color.rgb = BLUE
        textbox(s, 6.75, y + 0.35, 1.0, 0.6, n, size=19, color=WHITE if n == str(active) else BLUE, bold=True, align=PP_ALIGN.CENTER)
        textbox(s, 8.1, y + 0.35, 9.6, 0.6, t, size=16, color=WHITE if n == str(active) else BLUE, bold=True)
        y += 2.3
    return s


def meta(page, title, role, source, visual, claim):
    slides_meta.append((page, title, role, source, visual, claim))


page = 1
s = blank()
textbox(s, 0.0, 6.0, 25.4, 2.1, "基于规划反馈的图推理多智能体\n协同求解方法研究", size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 6.7, 10.8, 11.0, 1.8, "硕士研究生：王腾龙\n指 导 教 师：林丽", size=17, color=BLACK, align=PP_ALIGN.CENTER)
textbox(s, 6.6, 2.9, 12.0, 0.9, "东南大学硕士学位论文答辩报告", size=16, color=BLUE, align=PP_ALIGN.CENTER)
add_img(s, "framework.png", 17.8, 11.5, 5.7, 3.4)
meta(page, "封面", "正式答辩开场", "论文题名与作者信息", "右下角框架图", "本文研究基于规划反馈的图推理多智能体协同求解方法。"); page += 1

s = section(page, "提纲", 0)
meta(page, "提纲", "说明答辩结构", "学长模板第2页", "四章式目录", "答辩按背景现状、目标内容、方法系统、实验总结推进。"); page += 1
s = section(page, "提纲", 1)
meta(page, "提纲", "章节过渡", "学长模板第3页", "目录高亮", "先说明为什么图推理需要新的求解框架。"); page += 1

for title, img, pts, claim in [
    ("研究背景", "intro.png", ["大语言模型具备复杂推理与程序生成能力", "图数据具有拓扑依赖、强约束和难线性化特征", "图推理要求模型同时完成结构理解、算法匹配与结果验证"], "图推理不是一般文本推理的简单延伸。"),
    ("研究背景", "graph_reasoning_challenge.png", ["长序列图输入容易造成拓扑关系遗漏", "复杂约束下容易出现答案幻觉", "程序可运行并不等价于结果满足任务目标"], "端到端生成在复杂图任务中不稳定。"),
]:
    s = blank(); header(s, title, page)
    add_img(s, img, 1.0, 3.3, 10.9, 8.1)
    bullets(s, 13.2, 4.0, 10.2, 5.8, pts, size=16)
    textbox(s, 2.1, 15.3, 21.5, 0.8, claim, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    meta(page, title, "背景问题", "论文第1章及开题PPT背景页", img, claim); page += 1

s = blank(); header(s, "研究背景", page)
for i, (name, pts) in enumerate([("路径规划", ["最短路径", "旅行商问题", "交通与物流优化"]), ("网络分析", ["连通性", "最大流", "依赖关系分析"]), ("组合优化", ["最大团", "点覆盖", "匹配与搜索"]) ]):
    x = 1.2 + i * 8.0
    tag(s, x, 4.0, name, w=4.0)
    box(s, x, 5.0, 6.2, 5.2)
    bullets(s, x + 0.5, 5.7, 5.1, 3.0, pts, size=15)
textbox(s, 2.2, 14.0, 21.0, 1.0, "可靠图推理的关键不是生成更长解释，而是生成满足图结构约束的可验证结果。", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
meta(page, "研究背景", "应用价值", "论文第1章", "三类应用模块", "图推理具有明确应用价值，可靠性直接影响决策质量。"); page += 1

s = blank(); header(s, "研究现状", page)
add_img(s, "llm_graph_reasoning_taxonomy.png", 1.0, 3.1, 12.4, 8.2)
bullets(s, 14.5, 3.8, 9.0, 5.8, ["通用推理增强：CoT、ToT、GoT 等", "图任务方法：GraphWiz、GraphTeam、GCoder 等", "整体趋势：由单次生成走向结构化协同与执行验证"], size=15)
meta(page, "研究现状", "研究脉络", "论文第2章 图 llm_graph_reasoning_taxonomy", "二维分类图", "已有研究开始重视图任务专门化和系统化协同。"); page += 1

s = blank(); header(s, "研究现状", page)
for i, (h, b) in enumerate([("不足点1", "端到端推理负担重，复杂输入下容易遗漏拓扑与边界条件"), ("不足点2", "知识调用缺乏中间组织，检索结果与求解步骤不一定匹配"), ("不足点3", "程序生成缺少执行反馈，错误修复容易变成无方向重试")]):
    y = 3.3 + i * 3.5
    tag(s, 1.2, y, h, w=3.2)
    box(s, 4.8, y, 18.8, 1.5, b)
textbox(s, 2.1, 14.5, 21.4, 0.9, "本文切入点：以显式规划作前向约束，以执行反馈作反向修正。", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
meta(page, "研究现状", "不足归纳", "论文第1、2章", "三条不足", "本文不是单纯增加智能体，而是让规划、知识和验证形成闭环。"); page += 1

s = section(page, "提纲", 2); meta(page, "提纲", "章节过渡", "学长模板第12页", "目录高亮", "转入研究目标与研究内容。"); page += 1

s = blank(); header(s, "研究目标", page)
for i, (lab, txt) in enumerate([("总体\n目标", "提升复杂图推理任务的稳定性、正确性与可解释性"), ("理论\n目标", "构建可规划、可验证、可修正的程序化图推理求解范式"), ("系统\n目标", "实现主控调度、多智能体协作、共享黑板与执行反馈机制")]):
    y = 3.3 + i * 3.6
    tag(s, 1.5, y, lab, w=3.1)
    box(s, 5.0, y - 0.1, 18.2, 1.7, txt)
meta(page, "研究目标", "目标定义", "论文第1章、第6章", "三目标条", "本文目标是把图推理从一次性生成变成受控协同求解。"); page += 1

s = blank(); header(s, "机遇", page)
tag(s, 1.3, 4.0, "复杂图推理", w=5.0); tag(s, 10.0, 4.0, "结构化规划", w=5.0, color=ORANGE); tag(s, 18.3, 4.0, "可执行求解", w=5.0)
textbox(s, 7.0, 4.15, 2.0, 0.6, "→", size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 15.8, 4.15, 2.0, 0.6, "→", size=24, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
bullets(s, 2.0, 7.0, 20.5, 4.0, ["机遇1：引入结构化规划层，降低自然语言图任务的求解难度", "机遇2：以规划结果引导知识检索，提高 API 与算法知识的匹配精度", "机遇3：利用程序执行反馈，将错误定位和修复纳入闭环"], size=17)
meta(page, "机遇", "研究机会", "论文第3章、第5章检索实验", "三段流程", "规划层能够连接复杂输入、知识调用和可执行求解。"); page += 1

for no, t, pts in [
    ("1", "如何将自然语言图任务结构化解耦？", ["抽取任务类型、图结构、约束集合和输出格式", "将任务目标转化为算法级伪代码规划", "为检索、编码和验证提供统一中间语义"]),
    ("2", "如何使外部知识与求解步骤准确匹配？", ["以规划伪代码和结构约束标签构造查询", "检索图算法 API 与实现注意事项", "通过质量评估触发片段精炼或补充检索"]),
    ("3", "如何发现并修复程序实现偏差？", ["执行测试检查图构建和输出格式", "约束检查覆盖路径、容量、时间窗口等条件", "错误反馈驱动局部修复、补充检索或重生成"]),
]:
    s = blank(); header(s, "挑战", page)
    tag(s, 1.4, 3.5, no, w=1.2, color=ORANGE)
    textbox(s, 3.0, 3.35, 20.0, 0.8, t, size=22, color=BLUE, bold=True)
    bullets(s, 3.1, 6.0, 18.5, 5.0, pts, size=17)
    meta(page, "挑战", "挑战与应对", "论文第3章方法设计", f"挑战{no}要点", t); page += 1

s = blank(); header(s, "研究内容", page)
add_img(s, "research_roadmap.png", 1.0, 3.0, 9.8, 7.3)
for i, (h, txt) in enumerate([("内容一", "设计规划反馈图推理协同求解方法"), ("内容二", "实现主控调度与共享黑板支撑的多智能体系统"), ("内容三", "在五个公开基准上进行系统实验验证")]):
    y = 3.2 + i * 3.0
    tag(s, 12.0, y, h, w=3.2)
    box(s, 15.5, y, 8.0, 1.25, txt)
meta(page, "研究内容", "研究内容总览", "论文第1章研究内容；图 research_roadmap", "路线图", "本文工作覆盖方法设计、系统实现和实验验证。"); page += 1

s = section(page, "提纲", 3); meta(page, "提纲", "章节过渡", "学长模板第20页", "目录高亮", "进入技术路线和系统实现。"); page += 1

s = blank(); header(s, "系统总体设计", page)
add_img(s, "framework.png", 1.0, 3.0, 16.0, 9.0)
bullets(s, 18.0, 3.6, 5.8, 6.0, ["模块一：规划生成", "模块二：规划引导检索", "模块三：验证推理闭环"], size=15)
meta(page, "系统总体设计", "总体框架", "论文第3章 图 framework_overview", "框架图", "PlanGraph 将求解过程组织为规划、检索、编码和验证反馈闭环。"); page += 1

for module, title2, subq, idea, img in [
    ("规划生成方法", "任务解析", "子问题1：如何从自然语言问题中恢复任务结构？", "将输入规整为任务类型、图标签、约束集合和输出格式，降低原始文本噪声。", None),
    ("规划生成方法", "伪代码规划", "子问题2：如何把任务目标转化为可执行求解步骤？", "生成算法级伪代码，显式表达图构建、算法选择、约束检查和输出格式化步骤。", None),
    ("知识检索方法", "规划引导检索", "子问题3：如何召回与当前求解步骤匹配的算法知识？", "使用规划伪代码与结构约束标签构造查询，从 3349 条图算法知识条目中召回 API 与实现知识。", None),
    ("知识检索方法", "质量评估与补充检索", "子问题4：如何降低无关知识进入编码阶段的概率？", "对候选知识进行相关性评估，必要时触发片段精炼、查询重写或补充检索。", None),
    ("验证推理方法", "执行验证", "子问题5：如何判断候选程序是否真正满足任务目标？", "通过执行测试、约束检查和格式校验识别图构建错误、API 误用和输出不一致。", "fig3_2.png"),
    ("验证推理方法", "反馈修复", "子问题6：如何避免无方向的重复生成？", "根据错误类型选择局部修复、补充检索、重生成或上游回退，并设置有界尝试次数。", "fig3_2.png"),
]:
    s = blank(); header(s, f"技术路线-{module}-{title2}", page)
    box(s, 0.9, 2.75, 22.7, 1.45, color=BLUE_LIGHT)
    textbox(s, 1.1, 3.1, 21.8, 0.6, subq, size=17, color=BLUE, bold=True)
    tag(s, 1.0, 5.1, "解决思路", w=3.8)
    bullets(s, 1.2, 6.3, 10.5, 4.0, [idea], size=16)
    tag(s, 1.0, 11.0, "难点分析", w=3.8, color=ORANGE)
    bullets(s, 1.2, 12.2, 10.5, 2.3, ["图任务同时包含结构、算法和约束，单轮生成难以稳定覆盖。"], size=15)
    if img:
        add_img(s, img, 12.6, 5.0, 10.5, 7.3)
    else:
        for i, lab in enumerate(["输入", "中间工件", "后续阶段"]):
            tag(s, 13.0 + i * 3.2, 7.0, lab, w=2.5, color=BLUE if i != 1 else ORANGE)
            if i < 2:
                textbox(s, 15.4 + i * 3.2, 7.12, 1.0, 0.5, "→", size=20, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    meta(page, f"技术路线-{module}-{title2}", "方法细化页", "论文第3章", img or "流程示意", idea); page += 1

s = blank(); header(s, "系统总体设计", page)
add_img(s, "SystemArchitecture.png", 0.8, 2.8, 17.2, 9.4)
bullets(s, 18.6, 3.5, 5.0, 5.6, ["输入与任务建模层", "主控调度层", "多智能体执行层", "共享状态与支撑资源层"], size=14)
meta(page, "系统总体设计", "系统架构", "论文第4章 图 system_architecture", "系统架构图", "PlanGraph 的协同依赖受控调度，而不是自由对话。"); page += 1

for t, data, colw, claim in [
    ("系统实现-主控调度", [["状态","输入","产出","失败恢复"],["PARSE","原始问题","结构化任务","重新解析"],["PLAN","任务工件","规划工件","重新规划"],["RETRIEVE","规划工件","知识工件","补充检索"],["VERIFY","候选程序","反馈工件","修复/重生成"]], [4.0,5.5,5.5,6.0], "状态化调度使错误定位和恢复路径可控。"),
    ("系统实现-运行配置", [["配置项","默认值","作用"],["检索返回条数","5","控制候选知识范围"],["最大局部修复轮数","3","控制反馈闭环深度"],["最大重生成轮数","1","限制失败后探索成本"],["单次执行超时","15秒","避免复杂搜索失控"],["随机种子","42","支撑实验复现"]], [7.0,4.0,10.0], "反馈闭环必须有界，才能兼顾准确率和工程成本。"),
]:
    s = blank(); header(s, t, page)
    table(s, 1.4, 3.3, data, colw, fs=11)
    textbox(s, 2.0, 14.6, 21.0, 0.8, claim, size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    meta(page, t, "系统机制", "论文第4章", "表格", claim); page += 1

s = section(page, "提纲", 4); meta(page, "提纲", "章节过渡", "学长模板第39页", "目录高亮", "进入实验验证和总结展望。"); page += 1

for t, data, claim in [
    ("数据集与评价指标", [["基准测试集","侧重点","代表任务"],["GraphArena","经典图问题与组合优化","连通性、最大团、旅行商"],["NLGraph","自然语言图算法推理","最短路径、最大流"],["GraphWiz","指令式图推理","哈密顿路径、子图匹配"],["LLM4DyG","动态图与时间约束","时序路径、动态邻居"],["GraphInstruct","图指令跟随","邻居查询、遍历"]], "实验覆盖经典图算法、组合优化、指令跟随和动态图推理。"),
    ("实验对比算法", [["类别","方法","说明"],["直接推理","GPT-4o","通用强基座模型"],["多智能体","GraphTeam","代表性协同求解方法"],["程序生成","GCoder","代表性代码化图推理方法"],["本文方法","PlanGraph","规划反馈多智能体协同求解"]], "实验固定数据划分、模型版本、执行环境和答案判定规则。"),
]:
    s = blank(); header(s, t, page)
    table(s, 1.0, 3.1, data, [5.0,6.0,11.0], fs=10)
    textbox(s, 2.0, 14.4, 21.0, 0.8, claim, size=17, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    meta(page, t, "实验设置", "论文第5章", "实验设置表", claim); page += 1

s = blank(); header(s, "方法性能评估", page)
table(s, 0.7, 2.8, [["基准","任务数","GPT-4o","GraphTeam","GCoder","PlanGraph","提升"],["GraphArena","10","49.08","95.14","94.27","97.74","+2.60"],["NLGraph","8","51.40","97.71","96.90","98.72","+1.01"],["GraphWiz","9","47.61","88.62","90.20","97.01","+6.81"],["LLM4DyG","9","58.04","96.35","93.78","96.11","-0.24"],["GraphInstruct","21","51.14","93.48","92.56","94.87","+1.39"],["平均","--","51.45","94.26","93.54","96.89","+2.31"]], [3.5,2.2,3.0,3.2,3.0,3.3,2.6], fs=9)
bullets(s, 1.6, 12.5, 20.5, 2.2, ["PlanGraph 在五个基准中四个取得最优，平均 EM 达到 96.89%。", "GraphWiz 提升 6.81 个百分点，说明高约束路径搜索任务收益明显。"], size=15)
meta(page, "方法性能评估", "总体结果", "论文第5章 表 overall_result", "总体结果表", "PlanGraph 平均 EM 为 96.89%，较最佳现有方法平均提升 2.31 个百分点。"); page += 1

for t, img, pts in [
    ("任务类别实验结果", "exp_category_accuracy.png", ["五类任务上均取得最高精确匹配率", "哈密顿类任务平均 EM 为 95.33%", "相比最佳现有方法提升 21.03 个百分点"]),
    ("复杂度分层实验结果", "exp_complexity_trend.png", ["非 NP 完全任务接近饱和", "NP 完全任务上 PlanGraph 提升更明显", "复杂任务更能体现规划与验证闭环价值"]),
    ("样本难度实验结果", "exp_sample_difficulty_trend.png", ["难度提高后各方法均下降", "PlanGraph 在困难样本下降幅度更小", "大图搜索空间下执行验证收益更明显"]),
    ("消融实验", "exp_ablation.png", ["去除规划模块后下降最明显", "去除验证模块后 GraphWiz 从 97.01% 降至 90.72%", "检索模块提升温和但稳定"]),
]:
    s = blank(); header(s, t, page)
    add_img(s, img, 1.0, 3.0, 13.2, 8.3)
    bullets(s, 15.4, 4.0, 7.8, 5.5, pts, size=15)
    meta(page, t, "实验结果", "论文第5章", img, pts[0]); page += 1

s = blank(); header(s, "鲁棒性分析", page)
table(s, 3.0, 3.2, [["扰动类型","GPT-4o","GraphTeam","PlanGraph"],["节点顺序打乱","51.6","92.4","96.8"],["边描述改写","48.9","91.1","95.7"],["输出格式切换","44.3","89.8","94.9"]], [7.0,4.5,4.5,4.5], fs=12)
textbox(s, 2.2, 10.7, 21.0, 1.0, "结构化任务表示降低了表面措辞和输出形式变化对后续规划的影响。", size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
meta(page, "鲁棒性分析", "鲁棒性结果", "论文第5章 表 robustness", "鲁棒性表", "PlanGraph 在三类输入扰动下仍保持高准确率。"); page += 1

s = blank(); header(s, "参数敏感性与工程开销", page)
table(s, 0.8, 2.9, [["参数(k,r)","总体EM","困难EM","平均修复","耗时(s)"],["(1,0)","70.00","55.00","0.00","13.80"],["(3,1)","74.17","63.75","0.36","15.20"],["(7,3)","73.75","66.25","0.88","17.60"],["(5,4)","71.25","65.00","1.14","20.90"]], [4.5,3.5,3.5,3.5,3.5], fs=10)
table(s, 14.3, 8.0, [["任务","耗时","生成","验证"],["最短路径","10.8","1.1","1.3"],["最大流","14.6","1.3","1.8"],["旅行商","24.3","1.8","2.7"]], [3.6,2.4,2.4,2.4], fs=9)
meta(page, "参数敏感性与工程开销", "补充实验", "论文第5章 表 param_main_result / efficiency", "双表", "反馈深度需要适度控制，组合优化任务开销更高。"); page += 1

s = blank(); header(s, "样例分析", page)
add_img(s, "case2.png", 1.0, 3.0, 11.0, 8.5)
bullets(s, 13.0, 3.5, 10.2, 5.8, ["样例：五节点带权无向图旅行商问题", "规划显式表达访问全部节点、回到起点和总代价最小化", "验证检查路径闭合、节点覆盖、边合法性和代价计算"], size=15)
textbox(s, 13.0, 11.4, 10.2, 1.0, "最终路径：A → B → E → D → C → A，总代价 14", size=16, color=BLUE, bold=True)
meta(page, "样例分析", "典型案例", "论文第5章 TSP 案例", "案例图", "规划反馈能把组合优化问题转化为可验证的程序化求解过程。"); page += 1

s = blank(); header(s, "误差分析", page)
table(s, 1.0, 3.1, [["错误类型","具体表现","主要影响阶段"],["任务识别偏差","判定任务与优化任务混淆","规划阶段"],["约束遗漏","忽略起点、回路、时间窗口等限制","规划/编码"],["API误用","函数语义相近但约束不匹配","检索/编码"],["边界情况失效","验证样例未覆盖特殊结构","验证阶段"],["输出格式错误","路径、集合或文本格式不符合脚本","答案生成"]], [5.0,12.0,5.0], fs=10)
meta(page, "误差分析", "不足分析", "论文第5章 表 error_taxonomy", "错误类型表", "后续优化应集中在规划粒度、检索重排序和验证样例质量。"); page += 1

s = blank(); header(s, "论文总结", page)
bullets(s, 1.4, 3.0, 22.5, 9.0, ["系统分析了大语言模型在图推理任务中的结构理解、约束建模、算法映射和程序执行困难。", "提出 PlanGraph 框架，将规划生成、知识检索、程序生成、执行验证和反馈修复组织为协同闭环。", "实现了主控调度、共享黑板、中间工件、异常恢复和任务适配等系统机制。", "在五个公开图推理基准上完成系统评估，平均 EM 达到 96.89%，较最佳现有方法平均提升 2.31 个百分点。"], size=16)
meta(page, "论文总结", "结论总结", "论文第6章", "总结列表", "本文完成了从方法到系统再到验证的闭环研究。"); page += 1

s = blank(); header(s, "论文创新点", page)
for i, (h, txt) in enumerate([("创新点1", "提出基于规划反馈的图推理协同框架，使规划前向约束检索与编码，反馈反向修正实现偏差。"), ("创新点2", "提出伪代码规划引导的知识检索机制，将原始问题查询转化为贴近算法步骤的结构化查询。"), ("创新点3", "构建执行验证与反馈推理闭环，通过测试、约束检查和有界修复提升可靠性。")]):
    y = 3.2 + i * 3.6
    tag(s, 1.3, y, h, w=4.0, color=ORANGE if i == 1 else BLUE)
    box(s, 5.7, y, 17.8, 1.55, txt)
meta(page, "论文创新点", "创新点展示", "论文第6章", "三条创新", "三项创新均对应具体方法模块和实验验证。"); page += 1

s = blank(); header(s, "展望", page)
for i, (h, txt) in enumerate([("规划表示", "增强复杂时间依赖、动态图演化和异构关系表达能力"), ("知识检索", "利用失败反馈驱动知识重排序和补充检索"), ("验证机制", "构造更具代表性的边界样例，提升约束覆盖度"), ("真实应用", "面向知识图谱问答、软件依赖分析和网络运维诊断等场景落地")]):
    y = 3.1 + i * 2.8
    tag(s, 1.3, y, h, w=4.0)
    box(s, 5.7, y, 17.8, 1.2, txt)
meta(page, "展望", "后续工作", "论文第6章展望", "四条展望", "后续重点是增强中间表示与约束检查质量。"); page += 1

s = blank()
textbox(s, 0.0, 6.3, 25.4, 2.2, "感谢各位评委老师！\n请大家提出宝贵意见！", size=32, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
meta(page, "致谢", "结束页", "学长模板第54页", "致谢页", "答辩结束，进入提问环节。")


prs.save(PPTX_OUT)

lines = ["# 二轮优化版大纲", "", "| 页码 | 页面标题 | 页面功能 | 使用材料来源 | 图表/图片 | 核心表达句 |", "|---:|---|---|---|---|---|"]
for row in slides_meta:
    lines.append("| {} | {} | {} | {} | {} | {} |".format(*[str(x).replace("|", "；") for x in row]))
OUTLINE_OUT.write_text("\n".join(lines), encoding="utf-8")

REVIEW_OUT.write_text("""# 二轮对比反省与优化说明

## 对比发现

- 学长模板是 4:3 页面比例，上一版误用 16:9，导致整体气质偏离模板。
- 学长模板主色高度统一，核心蓝色为 #02409A；上一版使用了较多橙、绿、浅蓝卡片，视觉更像重设计稿。
- 学长模板标题短，多数页面采用“研究背景 / 研究现状 / 技术路线-模块-子模块”的标题体系；上一版标题偏长。
- 学长方法页常用“子问题—难点分析—解决思路”的讲解结构；上一版更像论文提纲式摘要。
- 学长实验页重视“图表 + 结论句”，上一版已有结果但章节节奏偏压缩。

## 本轮优化

- 改为 4:3 页面比例，顶部蓝色标题栏、左侧短竖线、底部页码风格向学长模板靠拢。
- 统一主色为 #02409A，减少多色卡片，字体统一为微软雅黑。
- 将方法部分拆成 6 页技术路线细分页，对应规划生成、检索、验证反馈，采用“子问题/难点分析/解决思路”结构。
- 页数由 41 页调整为 42 页，更接近学长正式答辩节奏，同时仍删除不适配的冗余页面。
- 保留论文关键实验：总体结果、任务类别、复杂度、样本难度、消融、鲁棒性、参数开销、案例和误差分析。
""", encoding="utf-8")

print(PPTX_OUT)
print(OUTLINE_OUT)
print(REVIEW_OUT)
