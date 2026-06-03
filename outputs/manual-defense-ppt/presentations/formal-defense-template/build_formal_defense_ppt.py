from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Cm, Pt


ROOT = Path("/Users/wtl/Desktop/我的毕业论文 ")
WORK = ROOT / "outputs/manual-defense-ppt/presentations/formal-defense-template"
ASSETS = WORK / "assets"
OUT_DIR = WORK / "output"
OUT_DIR.mkdir(parents=True, exist_ok=True)

PPTX_OUT = OUT_DIR / "王腾龙-硕士论文正式答辩PPT-初稿.pptx"
OUTLINE_OUT = OUT_DIR / "王腾龙-硕士论文正式答辩PPT-大纲.md"
ANALYSIS_OUT = OUT_DIR / "学长模板分析与调整说明.md"

prs = Presentation()
prs.slide_width = Cm(33.867)
prs.slide_height = Cm(19.05)

W = prs.slide_width
H = prs.slide_height

BLUE = RGBColor(31, 78, 121)
BLUE_DARK = RGBColor(20, 55, 87)
BLUE_LIGHT = RGBColor(221, 235, 247)
ORANGE = RGBColor(237, 125, 49)
GREEN = RGBColor(112, 173, 71)
GRAY = RGBColor(89, 89, 89)
LIGHT = RGBColor(248, 250, 252)
BLACK = RGBColor(30, 30, 30)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(192, 80, 77)

FONT_CN = "Microsoft YaHei"


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_fill(shape, color, transparency=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.transparency = transparency


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    set_fill(bg, WHITE)
    bg.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Cm(0.32))
    set_fill(bar, BLUE)
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Cm(0.32), Cm(4.3), Cm(0.16))
    set_fill(accent, ORANGE)
    accent.line.fill.background()


def txbox(slide, x, y, w, h, text="", font=FONT_CN, size=22, color=BLACK, bold=False,
          align=PP_ALIGN.LEFT, valign=None):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0] if p.runs else p.add_run()
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    if valign:
        tf.vertical_anchor = valign
    return shape


def title(slide, text, page=None, chapter=""):
    add_bg(slide)
    txbox(slide, 1.35, 0.65, 24.5, 1.1, text, size=24, bold=True, color=BLUE_DARK)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(1.35), Cm(1.92), Cm(3.7), Cm(0.08)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ORANGE
    slide.shapes[-1].line.fill.background()
    if chapter:
        txbox(slide, 26.3, 0.75, 5.2, 0.75, chapter, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    if page is not None:
        txbox(slide, 30.7, 18.1, 1.8, 0.4, f"{page:02d}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def bullet_list(slide, x, y, w, h, items, size=17, color=BLACK, line_spacing=1.0):
    shape = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT_CN
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(7)
        p.line_spacing = line_spacing
    return shape


def section_slide(num, section, subtitle, page):
    slide = blank()
    set_fill(slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H), BLUE_DARK)
    slide.shapes[-1].line.fill.background()
    txbox(slide, 2.0, 3.9, 2.2, 1.2, f"{num}", size=34, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txbox(slide, 4.8, 3.65, 18.5, 1.45, section, size=31, bold=True, color=WHITE)
    txbox(slide, 4.9, 5.28, 20.5, 1.0, subtitle, size=18, color=BLUE_LIGHT)
    txbox(slide, 28.8, 17.8, 2.2, 0.4, f"{page:02d}", size=9, color=BLUE_LIGHT, align=PP_ALIGN.RIGHT)
    return slide


def card(slide, x, y, w, h, head, body, color=BLUE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    set_fill(box, RGBColor(250, 252, 255))
    box.line.color.rgb = RGBColor(205, 218, 232)
    txbox(slide, x + 0.35, y + 0.25, w - 0.7, 0.55, head, size=15, bold=True, color=color)
    bullet_list(slide, x + 0.45, y + 1.05, w - 0.7, h - 1.25, body, size=12, color=BLACK)
    return box


def add_image(slide, img, x, y, w, h):
    p = ASSETS / img
    if p.exists():
        slide.shapes.add_picture(str(p), Cm(x), Cm(y), width=Cm(w), height=Cm(h))
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
        set_fill(ph, RGBColor(245, 245, 245))
        ph.line.color.rgb = RGBColor(180, 180, 180)
        txbox(slide, x + 0.4, y + h / 2 - 0.3, w - 0.8, 0.6, f"需替换图片：{img}", size=13, color=RED, align=PP_ALIGN.CENTER)


def simple_table(slide, x, y, rows, cols, data, col_widths=None, font_size=10, header=True):
    table = slide.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(sum(col_widths) if col_widths else 25), Cm(0.85 * rows)).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Cm(cw)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Cm(0.08)
            cell.margin_right = Cm(0.08)
            cell.margin_top = Cm(0.04)
            cell.margin_bottom = Cm(0.04)
            fill = BLUE if header and r == 0 else (RGBColor(241, 247, 252) if r % 2 == 1 else WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_CN
                    run.font.size = Pt(font_size)
                    run.font.bold = header and r == 0
                    run.font.color.rgb = WHITE if header and r == 0 else BLACK
    return table


slides_meta = []


def add_meta(page, t, role, source, visual, claim):
    slides_meta.append((page, t, role, source, visual, claim))


def agenda(slide, active=None):
    title(slide, "提纲", page=2 if active is None else active + 2)
    items = [
        ("1", "研究背景、研究现状"),
        ("2", "研究目标、研究内容"),
        ("3", "技术路线、系统实现"),
        ("4", "实验验证、总结展望"),
    ]
    y = 4.0
    for idx, text in items:
        fill = BLUE if active == int(idx) else RGBColor(242, 246, 250)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(5.0), Cm(y), Cm(23.5), Cm(1.65))
        set_fill(shape, fill)
        shape.line.color.rgb = RGBColor(210, 220, 230)
        txbox(slide, 5.55, y + 0.34, 1.1, 0.7, idx, size=20, bold=True, color=WHITE if active == int(idx) else BLUE, align=PP_ALIGN.CENTER)
        txbox(slide, 7.0, y + 0.36, 18.8, 0.7, text, size=18, bold=True, color=WHITE if active == int(idx) else BLUE_DARK)
        y += 2.35


# 1 cover
s = blank()
set_fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H), WHITE)
s.shapes[-1].line.fill.background()
set_fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(33.867), Cm(3.0)), BLUE_DARK)
s.shapes[-1].line.fill.background()
txbox(s, 2.0, 4.0, 23.0, 1.7, "基于规划反馈的图推理多智能体协同求解方法研究", size=27, bold=True, color=BLUE_DARK)
txbox(s, 2.05, 5.9, 17.0, 0.85, "硕士学位论文正式答辩报告", size=20, color=ORANGE, bold=True)
add_image(s, "framework.png", 20.2, 7.7, 10.0, 5.8)
txbox(s, 2.1, 10.7, 12.0, 2.2, "硕士研究生：王腾龙\n指导教师：林丽\n专业方向：人工智能 / 图推理与智能计算", size=16, color=BLACK)
txbox(s, 2.1, 16.35, 12.0, 0.6, "东南大学硕士学位论文答辩", size=14, color=GRAY)
add_meta(1, "封面", "正式答辩开场", "main.tex 论文题名与作者信息", "框架图弱化作背景证据", "本文研究基于规划反馈的图推理多智能体协同求解方法。")

# 2 agenda
s = blank(); agenda(s)
add_meta(2, "提纲", "说明答辩结构", "学长模板第2页", "四章式提纲", "答辩按背景现状、目标内容、方法系统、实验总结推进。")

page = 3
s = blank(); agenda(s, 1); add_meta(page, "提纲：背景与现状", "章节过渡", "学长模板第3页", "高亮第1章", "先说明为什么图推理需要新的求解框架。"); page += 1

s = blank(); title(s, "研究背景：大语言模型具备复杂推理潜力", page, "1 研究背景")
bullet_list(s, 2.1, 3.0, 13.2, 5.5, [
    "大语言模型通过自回归生成机制处理文本序列",
    "在自然语言理解、复杂推理和程序生成中表现突出",
    "多智能体协作使任务理解、规划、检索和执行可分阶段完成",
])
card(s, 18.0, 3.1, 10.6, 5.3, "研究切入点", ["将大模型能力从“文本回答”推进到“可规划、可验证、可修正”的图任务求解过程"], ORANGE)
add_meta(page, "研究背景：大语言模型具备复杂推理潜力", "背景铺垫", "论文摘要、第1章；开题PPT第4-5页", "左右图文结构", "大模型为图推理自动化提供可能，但需要系统化组织。"); page += 1

s = blank(); title(s, "研究背景：图数据具有强结构与强约束特征", page, "1 研究背景")
add_image(s, "intro.png", 2.0, 3.0, 13.8, 9.1)
bullet_list(s, 17.2, 3.4, 12.2, 6.5, [
    "拓扑依赖：节点与边关系决定推理过程",
    "约束显式：路径、匹配、流、回路等条件必须满足",
    "难以线性化：长序列输入容易造成结构遗漏",
])
txbox(s, 17.2, 10.6, 11.6, 1.2, "图推理不是一般文本推理的简单延伸，而是结构理解与可执行求解的复合任务。", size=17, bold=True, color=BLUE_DARK)
add_meta(page, "研究背景：图数据具有强结构与强约束特征", "问题对象定义", "论文第1章；图 intro", "论文图件PPT化", "图任务需要同时处理拓扑、约束和算法映射。"); page += 1

s = blank(); title(s, "应用场景：图推理支撑多类结构化决策", page, "1 研究背景")
for i, (h, b, c) in enumerate([
    ("路径规划", ["最短路径", "旅行商问题", "交通与物流优化"], BLUE),
    ("网络分析", ["连通性", "最大流", "社区与依赖关系"], GREEN),
    ("组合优化", ["最大团", "点覆盖", "独立集与匹配"], ORANGE),
]):
    card(s, 2.2 + i * 10.1, 3.1, 8.8, 6.7, h, b, c)
txbox(s, 3.2, 11.6, 26.0, 1.2, "关键要求：不仅要“答出结果”，还要保证结果满足图结构约束与算法目标。", size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "应用场景：图推理支撑多类结构化决策", "意义说明", "论文第1章；开题PPT第7页", "三类应用卡片", "图推理具有明确应用价值，可靠性直接影响决策质量。"); page += 1

s = blank(); title(s, "核心问题：端到端推理容易产生结构性错误", page, "1 研究背景")
add_image(s, "graph_reasoning_challenge.png", 2.0, 2.7, 14.6, 9.6)
bullet_list(s, 18.0, 3.2, 11.0, 6.4, [
    "长序列图输入导致拓扑关系遗漏",
    "复杂约束下容易出现答案幻觉",
    "程序可运行不代表结果满足任务目标",
])
txbox(s, 18.0, 10.4, 11.2, 1.4, "需求：把隐式推理拆解为可检查的中间步骤，并引入执行反馈。", size=17, bold=True, color=ORANGE)
add_meta(page, "核心问题：端到端推理容易产生结构性错误", "问题引出", "论文第1章；图 graph_reasoning_challenge", "挑战图+结论", "复杂图任务中，直接生成不稳定是主要瓶颈。"); page += 1

s = blank(); title(s, "研究现状：方法从单次生成走向结构化协同", page, "1 研究现状")
add_image(s, "llm_graph_reasoning_taxonomy.png", 2.0, 2.7, 16.0, 10.1)
bullet_list(s, 19.2, 3.1, 10.8, 7.0, [
    "通用推理增强：CoT、ToT、GoT 等扩展思维路径",
    "图任务方法：GraphWiz、GraphTeam、GCoder 等面向图算法求解",
    "发展趋势：从弱闭环走向系统化协同与可执行验证",
])
add_meta(page, "研究现状：方法从单次生成走向结构化协同", "国内外现状归纳", "论文第2章；图 llm_graph_reasoning_taxonomy", "二维研究脉络图", "已有研究已开始重视图任务专门化和系统化协同。"); page += 1

s = blank(); title(s, "现有不足：知识、规划与验证尚未形成闭环", page, "1 研究现状")
card(s, 2.2, 3.0, 8.6, 6.5, "不足一：端到端负担重", ["任务理解、算法选择、约束满足一次完成", "复杂输入下容易遗漏拓扑或边界条件"], BLUE)
card(s, 12.1, 3.0, 8.6, 6.5, "不足二：知识调用缺少组织", ["直接检索与当前求解步骤不一定匹配", "API 选择和约束解释容易偏差"], ORANGE)
card(s, 22.0, 3.0, 8.6, 6.5, "不足三：验证反馈不足", ["程序生成后缺少约束检查", "错误修复容易变成无方向重试"], RED)
txbox(s, 3.0, 11.4, 27.0, 1.1, "本文切入点：以显式规划作为前向约束，以执行反馈作为反向修正，构建图推理求解闭环。", size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "现有不足：知识、规划与验证尚未形成闭环", "不足归纳并引出本文", "论文第1、2章", "三栏不足卡片", "本文不是单纯增加智能体，而是组织规划、检索与验证反馈。"); page += 1

s = section_slide("2", "研究目标、研究内容", "从问题定位转入本文具体完成的工作", page)
add_meta(page, "章节页：研究目标、研究内容", "章节过渡", "学长模板第12页", "深蓝章节页", "接下来说明本文解决什么问题、完成哪些研究内容。"); page += 1

s = blank(); title(s, "研究目标：构建可规划、可验证、可修正的图推理框架", page, "2 目标内容")
for i, (h, b, c) in enumerate([
    ("总体目标", ["提升复杂图推理任务的稳定性、正确性与可解释性"], BLUE),
    ("理论目标", ["将自然语言图任务转化为结构化规划和可执行求解链路"], ORANGE),
    ("系统目标", ["实现多智能体协同、共享黑板、执行验证与反馈恢复机制"], GREEN),
]):
    card(s, 2.1 + i * 10.2, 3.0, 8.9, 6.4, h, b, c)
txbox(s, 3.2, 11.1, 26.2, 1.1, "目标对应论文第3章方法设计、第4章系统实现、第5章实验验证。", size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_meta(page, "研究目标：构建可规划、可验证、可修正的图推理框架", "目标定义", "论文第1章研究内容、第6章总结", "三目标卡片", "本文目标是把图推理从一次性生成变成受控协同求解。"); page += 1

s = blank(); title(s, "机遇一：结构化规划降低复杂任务求解负担", page, "2 目标内容")
txbox(s, 2.4, 3.2, 10.4, 1.0, "自然语言图任务", size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
txbox(s, 12.5, 3.2, 8.6, 1.0, "伪代码规划", size=17, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s, 22.0, 3.2, 8.0, 1.0, "可执行程序", size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
for x, c in [(2.2, BLUE_LIGHT), (12.0, RGBColor(255, 242, 229)), (21.6, RGBColor(235, 247, 233))]:
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(5.0), Cm(8.0), Cm(3.7)); set_fill(shp, c); shp.line.color.rgb = RGBColor(210, 220, 230)
bullet_list(s, 2.9, 5.55, 6.7, 2.5, ["识别任务类型", "抽取图结构", "整理约束"], size=14)
bullet_list(s, 12.7, 5.55, 6.7, 2.5, ["明确算法步骤", "显式表达目标", "连接知识检索"], size=14)
bullet_list(s, 22.3, 5.55, 6.7, 2.5, ["调用图算法", "执行与验证", "输出答案"], size=14)
for x in [10.6, 20.3]:
    txbox(s, x, 6.18, 1.0, 0.5, "→", size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txbox(s, 3.0, 11.4, 26.0, 1.0, "规划层的价值：把模型难以一次完成的复合任务拆解为可消费的中间语义。", size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "机遇一：结构化规划降低复杂任务求解负担", "机会分析", "论文第3章3.1-3.2", "三阶段流程草图", "显式规划能够降低任务理解、算法匹配和约束满足的同步负担。"); page += 1

s = blank(); title(s, "挑战一：如何将自然语言图任务结构化解耦？", page, "2 目标内容")
card(s, 2.2, 3.0, 8.9, 6.8, "输入复杂", ["图结构、任务目标和输出格式混在自然语言中", "节点、边、权重、方向等信息易遗漏"], BLUE)
card(s, 12.1, 3.0, 8.9, 6.8, "约束复杂", ["路径闭合、节点唯一访问、容量限制、时间窗口等条件并存", "单轮生成难以完整覆盖"], ORANGE)
card(s, 22.0, 3.0, 8.9, 6.8, "解法复杂", ["不同任务需要映射到不同算法族", "算法选择错误会导致后续全部偏移"], RED)
add_meta(page, "挑战一：如何将自然语言图任务结构化解耦？", "挑战定义", "论文第1、3章", "三类挑战卡片", "第一类挑战是从原始问题恢复可规划的任务表示。"); page += 1

s = blank(); title(s, "机遇二：规划引导检索可提升知识匹配精度", page, "2 目标内容")
simple_table(s, 2.1, 3.0, 5, 5, [
    ["查询方式", "Hit@1", "Hit@3", "Hit@5", "MRR"],
    ["原始问题文本", "42.6%", "61.8%", "72.9%", "51.4%"],
    ["任务类型+关键词", "68.9%", "84.7%", "92.1%", "76.3%"],
    ["规划伪代码", "76.4%", "90.8%", "96.3%", "83.5%"],
    ["规划+结构约束", "82.1%", "94.2%", "98.0%", "88.4%"],
], [6.0, 4.2, 4.2, 4.2, 4.2], font_size=11)
bullet_list(s, 24.2, 3.4, 7.2, 5.8, [
    "规划压缩了原始问题噪声",
    "结构约束标签提升排序质量",
    "知识检索从“相关文本”转向“求解步骤依据”",
], size=14)
add_meta(page, "机遇二：规划引导检索可提升知识匹配精度", "机制证据前置", "论文第5章 表 retrieval_quality", "检索质量表", "中间规划不仅服务编码，也改善知识调用质量。"); page += 1

s = blank(); title(s, "挑战二：如何避免知识调用和程序生成偏差？", page, "2 目标内容")
bullet_list(s, 2.2, 3.2, 13.0, 7.5, [
    "相似图算法 API 容易被误选或误用",
    "边权、方向、容量、时间等实现细节容易丢失",
    "程序能够运行但未必满足题目约束",
    "无验证的修复可能造成重复生成和错误强化",
])
add_image(s, "fig3_2.png", 17.0, 3.0, 12.0, 7.6)
add_meta(page, "挑战二：如何避免知识调用和程序生成偏差？", "挑战定义", "论文第3章、图 fig3_2", "反馈闭环图", "第二类挑战是让外部知识、程序实现和任务约束保持一致。"); page += 1

s = blank(); title(s, "研究内容：围绕 PlanGraph 闭环展开三项工作", page, "2 目标内容")
add_image(s, "research_roadmap.png", 2.0, 3.0, 12.6, 7.8)
card(s, 16.0, 2.8, 12.2, 2.45, "内容一：规划反馈方法", ["任务解析、伪代码规划、规划引导检索、执行验证与反馈修复"], BLUE)
card(s, 16.0, 5.8, 12.2, 2.45, "内容二：多智能体系统实现", ["主控调度、共享黑板、中间工件、异常恢复与任务适配"], ORANGE)
card(s, 16.0, 8.8, 12.2, 2.45, "内容三：系统实验验证", ["五个公开基准、分层分析、检索验证、鲁棒性、消融和案例"], GREEN)
add_meta(page, "研究内容：围绕 PlanGraph 闭环展开三项工作", "研究内容总览", "论文第1章研究内容；图 research_roadmap", "路线图+三项内容", "本文工作覆盖方法设计、系统实现与实证验证。"); page += 1

s = section_slide("3", "技术路线、系统实现", "说明 PlanGraph 如何工作以及如何落地", page)
add_meta(page, "章节页：技术路线、系统实现", "章节过渡", "学长模板第20页", "深蓝章节页", "进入方法和系统主体。"); page += 1

s = blank(); title(s, "总体框架：规划、检索、验证构成反馈闭环", page, "3 技术路线")
add_image(s, "framework.png", 1.7, 2.6, 20.0, 10.0)
bullet_list(s, 23.0, 3.0, 8.0, 6.8, [
    "问题智能体：结构化解析任务",
    "规划智能体：生成算法级伪代码",
    "检索智能体：召回图算法知识",
    "编码智能体：生成候选程序",
    "验证反馈：测试、修复与重生成",
], size=13)
add_meta(page, "总体框架：规划、检索、验证构成反馈闭环", "核心方法框架", "论文第3章 图 framework_overview", "总体框架图", "PlanGraph 将求解过程组织为可观测的闭环链路。"); page += 1

s = blank(); title(s, "任务解析与规划生成：把问题转为算法级伪代码", page, "3 技术路线")
simple_table(s, 2.0, 3.0, 4, 3, [
    ["阶段", "输入", "输出"],
    ["任务解析", "自然语言问题 + 图结构", "任务类型、图标签、约束集合、输出格式"],
    ["规划生成", "结构化任务表示", "算法步骤、关键约束、候选算法族"],
    ["规划修正", "验证或检索反馈", "细化后的规划工件"],
], [5.0, 9.0, 13.0], font_size=11)
txbox(s, 3.0, 9.5, 26.2, 1.1, "规划不是答案，而是后续检索、编码和验证共同使用的中间锚点。", size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "任务解析与规划生成：把问题转为算法级伪代码", "方法模块说明", "论文第3章3.2", "阶段输入输出表", "规划模块降低一次性生成负担，并明确后续阶段的消费对象。"); page += 1

s = blank(); title(s, "规划引导检索：从算法步骤召回实现知识", page, "3 技术路线")
bullet_list(s, 2.0, 3.2, 12.4, 7.2, [
    "以规划伪代码、任务类型和结构约束构造查询",
    "从 3349 条图算法知识条目中召回 API 与实现注意事项",
    "通过质量评估触发片段精炼或补充检索",
    "将知识工件写入共享黑板供编码阶段使用",
], size=15)
card(s, 17.0, 3.0, 10.8, 6.4, "设计要点", ["规划约束检索目标", "检索补充实现细节", "质量评估降低无关知识进入编码阶段的概率"], ORANGE)
add_meta(page, "规划引导检索：从算法步骤召回实现知识", "方法模块说明", "论文第3章3.3、第5章检索质量实验", "要点卡片", "检索模块的重点是让外部知识与当前求解步骤对齐。"); page += 1

s = blank(); title(s, "验证推理闭环：用执行结果反向修正实现偏差", page, "3 技术路线")
add_image(s, "fig3_2.png", 2.0, 2.8, 14.3, 8.8)
bullet_list(s, 18.0, 3.1, 11.4, 6.2, [
    "候选程序必须通过执行测试和约束检查",
    "错误反馈驱动局部修复、补充检索或重生成",
    "有界重试避免无方向反复生成",
    "最终输出经过格式化与答案一致性检查",
], size=15)
add_meta(page, "验证推理闭环：用执行结果反向修正实现偏差", "方法模块说明", "论文第3章 图 artifact_feedback / fig3_2", "反馈闭环图", "执行反馈使系统能够发现并修复程序表面可运行但结果错误的问题。"); page += 1

s = blank(); title(s, "系统总体架构：主控调度器组织多智能体协同", page, "3 系统实现")
add_image(s, "SystemArchitecture.png", 1.6, 2.5, 21.0, 10.6)
bullet_list(s, 23.8, 3.0, 7.2, 6.0, [
    "输入与任务建模层",
    "主控调度层",
    "多智能体执行层",
    "共享状态与支撑资源层",
], size=14)
add_meta(page, "系统总体架构：主控调度器组织多智能体协同", "系统实现说明", "论文第4章 图 system_architecture", "系统架构图", "PlanGraph 的协同依赖受控调度，而不是自由对话。"); page += 1

s = blank(); title(s, "主控调度与共享黑板：保证过程可追踪、可恢复", page, "3 系统实现")
simple_table(s, 2.0, 2.9, 5, 4, [
    ["状态", "主要输入", "产出工件", "失败恢复"],
    ["PARSE", "原始问题", "结构化任务", "重新解析"],
    ["PLAN", "任务工件", "规划工件", "重新规划"],
    ["RETRIEVE", "规划工件", "知识工件", "补充检索"],
    ["VERIFY", "候选程序", "验证反馈", "局部修复/重生成"],
], [4.2, 7.0, 7.0, 7.0], font_size=10)
txbox(s, 3.0, 10.3, 26.0, 1.2, "共享黑板记录阶段事件和工件索引，使失败定位和实验复核建立在真实运行证据上。", size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "主控调度与共享黑板：保证过程可追踪、可恢复", "系统机制说明", "论文第4章 表 scheduler_states / blackboard_events", "状态转移表", "受控调度让多智能体协作具备边界和恢复路径。"); page += 1

s = blank(); title(s, "运行配置：用边界条件约束反馈成本", page, "3 系统实现")
simple_table(s, 3.0, 2.7, 8, 3, [
    ["配置项", "默认值", "作用"],
    ["检索返回条数", "5", "控制候选知识范围"],
    ["最大局部修复轮数", "3", "控制反馈闭环深度"],
    ["最大重生成轮数", "1", "限制失败后探索成本"],
    ["最大规划尝试次数", "2", "限制上游反复规划"],
    ["最大检索尝试次数", "2", "控制补充检索次数"],
    ["单次执行超时", "15 秒", "避免复杂搜索失控"],
    ["随机种子", "42", "支撑实验复现"],
], [7.0, 4.0, 14.5], font_size=10)
add_meta(page, "运行配置：用边界条件约束反馈成本", "系统配置说明", "论文第4章 表 runtime_config", "配置表", "反馈闭环必须有界，才能在准确率和成本之间取得平衡。"); page += 1

s = blank(); title(s, "任务适配：稳定主流程支持不同图任务扩展", page, "3 系统实现")
card(s, 2.2, 3.0, 8.8, 6.2, "稳定主流程", ["解析", "规划", "检索", "编码", "验证", "输出"], BLUE)
card(s, 12.2, 3.0, 8.8, 6.2, "可替换任务知识", ["图类型标签", "约束集合", "算法族", "规划模板", "验证规则"], ORANGE)
card(s, 22.2, 3.0, 8.8, 6.2, "扩展方式", ["补充任务标签", "注册验证规则", "扩展知识条目", "保持调度器不变"], GREEN)
txbox(s, 3.1, 11.0, 27.0, 1.0, "扩展性来自主流程、任务适配器和验证规则之间的职责分离。", size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "任务适配：稳定主流程支持不同图任务扩展", "系统扩展说明", "论文第4章任务适配策略", "三栏机制卡片", "PlanGraph 不为单一数据集写固定流程，而是通过适配器承载任务语义。"); page += 1

s = section_slide("4", "实验验证、总结展望", "用结果证明方法有效，并给出贡献与边界", page)
add_meta(page, "章节页：实验验证、总结展望", "章节过渡", "学长模板第39页", "深蓝章节页", "进入正式答辩重点：实验结果与结论。"); page += 1

s = blank(); title(s, "实验设置：五个公开图推理基准覆盖多类任务", page, "4 实验验证")
simple_table(s, 2.0, 2.8, 6, 3, [
    ["基准测试集", "侧重点", "代表任务"],
    ["GraphArena", "经典图问题与组合优化", "连通性、最大团、旅行商"],
    ["NLGraph", "自然语言图算法推理", "最短路径、最大流、拓扑排序"],
    ["GraphWiz", "指令式图推理", "哈密顿路径、子图匹配"],
    ["LLM4DyG", "动态图与时间约束", "时序路径、动态邻居查询"],
    ["GraphInstruct", "图指令跟随", "邻居查询、遍历、连通性"],
], [5.0, 10.0, 11.5], font_size=10)
txbox(s, 3.0, 11.1, 26.0, 0.8, "主指标：精确匹配率（Exact Match, EM）；失败、超时、格式错误均计入统计。", size=15, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
add_meta(page, "实验设置：五个公开图推理基准覆盖多类任务", "实验设计", "论文第5章 数据集与评价指标", "数据集表", "实验覆盖经典图算法、组合优化、指令跟随和动态图推理。"); page += 1

s = blank(); title(s, "对比方法与实现设置：固定口径评估方法本身", page, "4 实验验证")
card(s, 2.0, 3.0, 7.0, 6.5, "对比方法", ["GPT-4o 直接推理", "GraphTeam", "GCoder", "公开最优结果"], BLUE)
card(s, 10.2, 3.0, 8.8, 6.5, "实现环境", ["GPT-4o 基座", "Python + NetworkX", "温度 0.2", "最大输出 2048 tokens"], ORANGE)
card(s, 20.4, 3.0, 9.2, 6.5, "公平设置", ["固定数据划分", "固定样本顺序", "固定模型版本", "统一答案判定规则"], GREEN)
add_meta(page, "对比方法与实现设置：固定口径评估方法本身", "实验可信度说明", "论文第5章 对比方法与实现设置", "三栏设置卡片", "实验尽量固定外部条件，比较规划、检索和验证模块带来的增益。"); page += 1

s = blank(); title(s, "总体结果：PlanGraph 在四个基准上取得最优", page, "4 实验验证")
simple_table(s, 1.5, 2.6, 7, 7, [
    ["基准", "任务数", "GPT-4o", "GraphTeam", "GCoder", "PlanGraph", "提升"],
    ["GraphArena", "10", "49.08", "95.14", "94.27", "97.74", "+2.60"],
    ["NLGraph", "8", "51.40", "97.71", "96.90", "98.72", "+1.01"],
    ["GraphWiz", "9", "47.61", "88.62", "90.20", "97.01", "+6.81"],
    ["LLM4DyG", "9", "58.04", "96.35", "93.78", "96.11", "-0.24"],
    ["GraphInstruct", "21", "51.14", "93.48", "92.56", "94.87", "+1.39"],
    ["平均", "--", "51.45", "94.26", "93.54", "96.89", "+2.31"],
], [4.0, 2.6, 3.6, 3.8, 3.5, 3.9, 3.0], font_size=9)
bullet_list(s, 3.0, 10.2, 25.5, 2.0, [
    "GraphWiz 提升最明显，说明规划反馈对算法约束强、路径搜索复杂的任务尤为有效。",
    "LLM4DyG 略低于 GraphTeam，提示动态图细粒度时序规划仍有改进空间。",
], size=14)
add_meta(page, "总体结果：PlanGraph 在四个基准上取得最优", "核心实验结果", "论文第5章 表 overall_result", "总体结果表", "PlanGraph 平均 EM 为 96.89%，较最佳现有方法平均提升 2.31 个百分点。"); page += 1

s = blank(); title(s, "任务类别分析：哈密顿类任务收益最明显", page, "4 实验验证")
add_image(s, "exp_category_accuracy.png", 2.0, 2.8, 16.2, 8.8)
bullet_list(s, 19.5, 3.0, 10.5, 6.2, [
    "五类任务上均取得最高精确匹配率",
    "哈密顿类任务平均 EM 95.33%",
    "相比最佳现有方法提升 21.03 个百分点",
    "优势来自约束显式化与可执行验证",
], size=14)
add_meta(page, "任务类别分析：哈密顿类任务收益最明显", "分层结果", "论文第5章 图 category_accuracy_comparison", "类别准确率图", "PlanGraph 的优势集中体现于高约束、强组合性的图任务。"); page += 1

s = blank(); title(s, "复杂度分析：NP 完全与困难样本更能体现闭环价值", page, "4 实验验证")
add_image(s, "exp_complexity_trend.png", 2.0, 3.0, 13.8, 8.4)
add_image(s, "exp_sample_difficulty_trend.png", 17.0, 3.0, 12.0, 8.4)
txbox(s, 3.0, 11.6, 26.0, 1.1, "结论：简单任务差距较小；进入 NP 完全或大图困难样本后，规划、检索与验证的协同收益更加明显。", size=16, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "复杂度分析：NP 完全与困难样本更能体现闭环价值", "分层结果", "论文第5章 图 np_complete_performance / sample_difficulty_performance", "复杂度与难度趋势图", "方法收益主要来自复杂任务，而非简单任务的边际提升。"); page += 1

s = blank(); title(s, "稳定性与检索质量：规划同时支撑检索和恢复", page, "4 实验验证")
simple_table(s, 1.8, 2.9, 4, 4, [
    ["数据集", "EM 波动", "首次通过率", "修复后收敛率"],
    ["GraphArena", "±0.7", "79.6", "91.4"],
    ["NLGraph", "±0.5", "82.3", "93.1"],
    ["GraphWiz", "±0.9", "74.8", "89.7"],
], [4.7, 4.3, 4.5, 5.0], font_size=10)
simple_table(s, 19.0, 2.9, 5, 3, [
    ["查询方式", "Hit@5", "MRR"],
    ["原始问题", "72.9%", "51.4%"],
    ["任务+关键词", "92.1%", "76.3%"],
    ["规划伪代码", "96.3%", "83.5%"],
    ["规划+约束", "98.0%", "88.4%"],
], [5.2, 3.6, 3.6], font_size=10)
txbox(s, 3.0, 11.0, 26.0, 1.0, "规划既是检索查询的结构化表达，也是验证修复时定位错误的中间锚点。", size=16, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "稳定性与检索质量：规划同时支撑检索和恢复", "机制验证", "论文第5章 表 stability_observation / retrieval_quality", "双表并排", "稳定运行和高质量检索共同支撑最终性能。"); page += 1

s = blank(); title(s, "参数与工程开销：反馈深度需要适度控制", page, "4 实验验证")
simple_table(s, 1.4, 2.6, 5, 5, [
    ["参数(k,r)", "总体EM", "困难EM", "平均修复", "耗时(s)"],
    ["(1,0)", "70.00", "55.00", "0.00", "13.80"],
    ["(3,1)", "74.17", "63.75", "0.36", "15.20"],
    ["(7,3)", "73.75", "66.25", "0.88", "17.60"],
    ["(5,4)", "71.25", "65.00", "1.14", "20.90"],
], [4.6, 3.8, 3.8, 4.0, 4.0], font_size=10)
simple_table(s, 20.6, 2.6, 4, 4, [
    ["任务", "耗时(s)", "生成轮数", "验证次数"],
    ["最短路径", "10.8", "1.1", "1.3"],
    ["最大流", "14.6", "1.3", "1.8"],
    ["旅行商", "24.3", "1.8", "2.7"],
], [4.0, 3.0, 3.0, 3.0], font_size=10)
bullet_list(s, 3.0, 10.2, 25.0, 2.0, [
    "无反馈修复时困难样本准确率明显下降；继续增大检索条数和修复轮数收益趋于饱和。",
    "组合优化任务开销更高，但换来更可信的约束满足与结果验证。",
], size=13)
add_meta(page, "参数与工程开销：反馈深度需要适度控制", "补充实验", "论文第5章 表 param_main_result / efficiency", "参数表+开销表", "PlanGraph 是成本换可靠性的设计，反馈闭环必须有界。"); page += 1

s = blank(); title(s, "鲁棒性分析：输入扰动下仍保持高准确率", page, "4 实验验证")
simple_table(s, 5.0, 3.0, 4, 4, [
    ["扰动类型", "GPT-4o", "GraphTeam", "PlanGraph"],
    ["节点顺序打乱", "51.6", "92.4", "96.8"],
    ["边描述改写", "48.9", "91.1", "95.7"],
    ["输出格式切换", "44.3", "89.8", "94.9"],
], [7.0, 5.0, 5.0, 5.0], font_size=12)
txbox(s, 4.2, 9.3, 24.5, 1.5, "原因：问题智能体先规整输入表达，程序执行结果在最终答案生成前还会经过格式化处理。", size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "鲁棒性分析：输入扰动下仍保持高准确率", "补充实验", "论文第5章 表 robustness", "鲁棒性表", "结构化任务表示降低了表面措辞和输出形式变化的影响。"); page += 1

s = blank(); title(s, "消融实验：规划模块、检索模块、验证模块均有贡献", page, "4 实验验证")
add_image(s, "exp_ablation.png", 2.0, 2.8, 16.0, 8.8)
bullet_list(s, 19.2, 3.2, 10.6, 6.2, [
    "去除规划模块后下降最明显",
    "去除验证模块后 GraphWiz 从 97.01% 降至 90.72%",
    "检索模块提升相对温和但稳定",
    "三者共同形成完整求解闭环",
], size=14)
add_meta(page, "消融实验：规划模块、检索模块、验证模块均有贡献", "关键机制验证", "论文第5章 图 ablation_result_visual", "消融实验图", "总体增益不是单点技巧，而是规划、知识和验证共同作用。"); page += 1

s = blank(); title(s, "案例分析：旅行商问题中的规划反馈求解过程", page, "4 实验验证")
add_image(s, "case2.png", 1.8, 2.7, 15.0, 9.3)
bullet_list(s, 18.0, 3.1, 11.6, 6.4, [
    "样例：五节点带权无向图，从 A 出发访问所有节点并返回起点",
    "规划显式表达哈密顿回路约束和总代价最小化目标",
    "验证检查节点覆盖、路径闭合、边合法性和代价计算",
])
txbox(s, 18.0, 10.0, 11.4, 1.1, "最终路径：A → B → E → D → C → A，总代价 14", size=16, bold=True, color=ORANGE)
add_meta(page, "案例分析：旅行商问题中的规划反馈求解过程", "典型案例", "论文第5章 图 tsp_case_pipeline / 案例描述", "案例流程图", "PlanGraph 能把组合优化问题转化为可验证的程序化求解过程。"); page += 1

s = blank(); title(s, "误差分析：当前方法的主要边界", page, "4 总结展望")
simple_table(s, 2.0, 2.8, 6, 3, [
    ["错误类型", "具体表现", "主要影响阶段"],
    ["任务识别偏差", "判定任务与优化任务混淆", "规划阶段"],
    ["约束遗漏", "忽略起点、回路、时间窗口等限制", "规划/编码"],
    ["API 误用", "函数语义相近但约束不匹配", "检索/编码"],
    ["边界情况失效", "验证样例未覆盖特殊结构", "验证阶段"],
    ["输出格式错误", "路径、集合或文本格式不符合脚本", "答案生成"],
], [5.0, 12.0, 8.0], font_size=10)
add_meta(page, "误差分析：当前方法的主要边界", "不足说明", "论文第5章 表 error_taxonomy", "错误类型表", "PlanGraph 的后续优化方向集中在规划粒度、检索重排序和验证样例质量。"); page += 1

s = blank(); title(s, "论文总结：本文完成了从方法到系统再到验证的闭环研究", page, "4 总结展望")
bullet_list(s, 2.0, 3.0, 27.5, 7.5, [
    "系统分析了大语言模型在图推理中的结构理解、约束建模、算法映射和程序执行困难。",
    "提出 PlanGraph 框架，将规划生成、知识检索、程序生成、执行验证和反馈修复组织为协同闭环。",
    "实现了主控调度、共享黑板、中间工件、异常恢复和任务适配等系统机制。",
    "在五个公开图推理基准上完成系统评估，平均 EM 达到 96.89%，较最佳现有方法平均提升 2.31 个百分点。",
])
add_meta(page, "论文总结：本文完成了从方法到系统再到验证的闭环研究", "结论总结", "论文第6章论文总结", "要点列表", "本文证明规划反馈协同机制能够提升复杂图推理的稳定性与可验证性。"); page += 1

s = blank(); title(s, "主要创新点：规划反馈驱动的图推理协同求解", page, "4 总结展望")
card(s, 2.0, 3.0, 9.0, 6.7, "创新点一", ["提出基于规划反馈的图推理协同框架", "使规划前向约束检索与编码，反馈反向修正实现偏差"], BLUE)
card(s, 12.2, 3.0, 9.0, 6.7, "创新点二", ["提出伪代码规划引导的知识检索机制", "将原始自然语言查询转化为更贴近算法步骤的结构化查询"], ORANGE)
card(s, 22.4, 3.0, 9.0, 6.7, "创新点三", ["构建执行验证与反馈推理闭环", "通过测试、约束检查和有界修复提升程序化图推理可靠性"], GREEN)
add_meta(page, "主要创新点：规划反馈驱动的图推理协同求解", "创新点展示", "论文第6章主要工作与结论", "三栏创新点", "三项创新均对应具体方法模块和实验验证。"); page += 1

s = blank(); title(s, "不足与展望：面向更复杂图任务继续增强中间表示", page, "4 总结展望")
card(s, 2.0, 3.0, 8.8, 6.5, "规划表示", ["增强对复杂时间依赖、动态图演化和异构关系的表达能力"], BLUE)
card(s, 12.0, 3.0, 8.8, 6.5, "知识检索", ["利用失败反馈驱动知识重排序和补充检索"], ORANGE)
card(s, 22.0, 3.0, 8.8, 6.5, "验证机制", ["构造更具代表性的边界样例，提升约束覆盖度"], GREEN)
txbox(s, 3.0, 11.0, 26.0, 1.0, "展望重点是增强中间表示与约束检查质量，而不是简单增加生成次数。", size=17, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_meta(page, "不足与展望：面向更复杂图任务继续增强中间表示", "展望收束", "论文第6章论文展望", "三栏展望卡片", "当前框架已具备稳定性，后续应提升动态图和长尾约束场景能力。"); page += 1

s = blank()
set_fill(s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H), BLUE_DARK)
s.shapes[-1].line.fill.background()
txbox(s, 4.0, 6.0, 25.0, 1.4, "感谢各位评委老师！", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, 4.0, 8.4, 25.0, 0.9, "请大家提出宝贵意见", size=22, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
add_meta(page, "致谢", "结束页", "学长模板第54页", "深蓝致谢页", "答辩结束，进入提问环节。")


# Write analysis and outline.
analysis = """# 学长模板分析与调整说明

## 1. 学长 PPT 总体结构

- 总页数：54 页。
- 章节结构：封面、提纲、研究背景与研究现状、研究目标与研究内容、技术路线与系统实现、实验验证与总结展望、成果与致谢。
- 大致页数：背景现状约 8 页；目标内容约 8 页；技术路线和系统实现约 19 页；实验验证约 12 页；总结成果致谢约 3 页。
- 页面功能：封面和提纲页建立答辩框架；背景页用问题与应用价值铺垫；现有工作页用二维分类矩阵引出不足；目标内容页用“机遇/挑战/研究内容”组织研究切入点；方法页按模块逐层展开；实验页先给数据集和基线，再给主结果、消融、低资源/参数/案例等补充结果；最后总结论文工作和成果。
- 版式规律：16:9 宽屏，主色为深蓝，橙色强调；标题通常位于左上，页脚有页码；章节页和提纲页反复出现，形成答辩节奏；多数正文页为左图右文、三栏卡片或表格+结论句。
- 叙事节奏：背景铺垫适中，方法模块展开较充分，实验结果集中在后段，最后用总结和成果收束。

## 2. 对本文 PPT 的调整

- 保留：封面、提纲复现、四章式组织、章节过渡页、左图右文、三栏挑战/创新、表格+结论句等模板逻辑。
- 调整：学长模板中大量“商品知识推理/文本编码/知识融合”细分页不适用于本文，因此替换为 PlanGraph 的任务解析、规划生成、规划引导检索、验证反馈、系统调度和实验分析页。
- 页数：压缩为 41 页，保留正式答辩所需完整链条，同时避免 54 页模板中不适配内容造成冗余。
- 优化：所有残留的电商客服内容均替换为论文内容；实验结果页增加每页核心结论，突出总体结果、任务类别、复杂度、稳定性、鲁棒性、消融和案例。
"""
ANALYSIS_OUT.write_text(analysis, encoding="utf-8")

lines = [
    "# 王腾龙硕士论文正式答辩 PPT 大纲",
    "",
    "| 页码 | 页面标题 | 页面功能 | 使用的论文材料来源 | 建议图表/图片 | 核心表达句 |",
    "|---:|---|---|---|---|---|",
]
for row in slides_meta:
    lines.append("| {} | {} | {} | {} | {} | {} |".format(*[str(x).replace("|", "；") for x in row]))
OUTLINE_OUT.write_text("\n".join(lines), encoding="utf-8")

prs.save(PPTX_OUT)
print(PPTX_OUT)
print(OUTLINE_OUT)
print(ANALYSIS_OUT)
